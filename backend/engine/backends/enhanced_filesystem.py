"""
EnhancedFilesystemBackend - 扩展官方 FilesystemBackend

遵循 DeepAgent 官方指导：
- 继承 FilesystemBackend，只重写 read() 方法
- 支持 DOCX/PDF/XLSX 等二进制文档格式
- 保持与官方 API 完全兼容（offset, limit 参数）
- 使用 LangChain 官方 document_loaders
"""

import os
import json
from pathlib import Path
from typing import List

from deepagents.backends.filesystem import FilesystemBackend
from deepagents.backends.utils import format_content_with_line_numbers, check_empty_content

# #region agent log
_DEBUG_LOG_PATH = Path(__file__).resolve().parents[3] / ".cursor" / "debug.log"
_ENABLE_FS_DEBUG_LOG = os.environ.get("ENABLE_FS_DEBUG_LOG", "").lower() in ("1", "true", "yes")


def _debug_log(location: str, message: str, data: dict = None):
    """写入调试日志到 NDJSON 文件"""
    if not _ENABLE_FS_DEBUG_LOG:
        return
    import time
    entry = {"location": location, "message": message, "data": data or {}, "timestamp": int(time.time() * 1000)}
    try:
        with open(_DEBUG_LOG_PATH, "a") as f:
            f.write(json.dumps(entry, ensure_ascii=False) + "\n")
    except Exception:
        pass
# #endregion

# 支持的二进制文档格式
BINARY_DOCUMENT_FORMATS = {'.docx', '.doc', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt'}


class EnhancedFilesystemBackend(FilesystemBackend):
    """
    扩展 FilesystemBackend 以支持二进制文档格式
    
    扩展功能：
    - DOCX/DOC: 使用 UnstructuredWordDocumentLoader
    - PDF: 使用 PDFPlumberLoader
    - XLSX/XLS: 使用 UnstructuredExcelLoader
    
    保持兼容：
    - 完全兼容官方 read() API（file_path, offset, limit）
    - 文本文件仍使用父类方法
    """
    
    def read(
        self,
        file_path: str,
        offset: int = 0,
        limit: int = 2000,
    ) -> str:
        """
        读取文件内容，支持二进制文档格式
        
        Args:
            file_path: 文件路径（绝对或相对）
            offset: 起始行偏移（0-indexed）
            limit: 最大读取行数
            
        Returns:
            格式化的文件内容（带行号），或错误消息
        """
        # #region agent log
        _debug_log("enhanced_filesystem.py:read:entry", "EnhancedFilesystemBackend.read 被调用", {
            "hypothesisId": "C",
            "file_path_arg": file_path,
            "backend_cwd": str(self.cwd),
            "virtual_mode": self.virtual_mode,
        })
        # #endregion
        resolved_path = self._resolve_path(file_path)
        # #region agent log
        _debug_log("enhanced_filesystem.py:read:resolved", "路径解析结果", {
            "hypothesisId": "C",
            "file_path_arg": file_path,
            "resolved_path": str(resolved_path),
            "exists": resolved_path.exists(),
            "is_file": resolved_path.is_file() if resolved_path.exists() else False,
        })
        # #endregion
        
        if not resolved_path.exists() or not resolved_path.is_file():
            return f"Error: File '{file_path}' not found"
        
        ext = resolved_path.suffix.lower()
        
        # 二进制文档使用专门的解析器
        if ext in BINARY_DOCUMENT_FORMATS:
            return self._read_binary_document(file_path, resolved_path, offset, limit)
        
        # 文本文件使用父类方法
        return super().read(file_path, offset, limit)
    
    async def aread(
        self,
        file_path: str,
        offset: int = 0,
        limit: int = 2000,
    ) -> str:
        """异步读取文件（使用 asyncio.to_thread）"""
        import asyncio
        return await asyncio.to_thread(self.read, file_path, offset, limit)
    
    def _read_binary_document(
        self,
        file_path: str,
        resolved_path: Path,
        offset: int,
        limit: int,
    ) -> str:
        """
        读取二进制文档（DOCX/PDF/XLSX）
        
        使用 LangChain 官方 document_loaders
        """
        ext = resolved_path.suffix.lower()
        
        try:
            lines = self._extract_document_lines(resolved_path, ext)
            
            if not lines:
                return check_empty_content("") or "File exists but has empty contents"
            
            # 应用 offset 和 limit
            start_idx = offset
            end_idx = min(start_idx + limit, len(lines))
            
            if start_idx >= len(lines):
                return f"Error: Line offset {offset} exceeds content length ({len(lines)} lines)"
            
            selected_lines = lines[start_idx:end_idx]
            return format_content_with_line_numbers(selected_lines, start_line=start_idx + 1)
            
        except ImportError as e:
            return f"Error: Missing dependency for {ext} files. Install: {e}"
        except Exception as e:
            return f"Error reading {ext} file '{file_path}': {e}"
    
    def _extract_document_lines(self, resolved_path: Path, ext: str) -> List[str]:
        """
        从文档中提取文本行
        
        使用 LangChain 官方 document_loaders
        """
        if ext in {'.docx', '.doc'}:
            return self._extract_word_lines(resolved_path)
        elif ext == '.pdf':
            return self._extract_pdf_lines(resolved_path)
        elif ext in {'.xlsx', '.xls'}:
            return self._extract_excel_lines(resolved_path)
        elif ext in {'.pptx', '.ppt'}:
            return self._extract_ppt_lines(resolved_path)
        else:
            raise ValueError(f"Unsupported format: {ext}")
    
    def _extract_word_lines(self, path: Path) -> List[str]:
        """提取 Word 文档内容"""
        try:
            from langchain_community.document_loaders import UnstructuredWordDocumentLoader
            loader = UnstructuredWordDocumentLoader(str(path))
            docs = loader.load()
            lines = []
            for doc in docs:
                lines.extend(doc.page_content.splitlines())
            return lines
        except ImportError:
            # 备选方案：使用 python-docx
            try:
                from docx import Document
                doc = Document(path)
                return [p.text for p in doc.paragraphs if p.text.strip()]
            except ImportError:
                raise ImportError("pip install python-docx or unstructured")
    
    def _extract_pdf_lines(self, path: Path) -> List[str]:
        """提取 PDF 文档内容"""
        try:
            from langchain_community.document_loaders import PDFPlumberLoader
            loader = PDFPlumberLoader(str(path))
            docs = loader.load()
            lines = []
            for doc in docs:
                lines.extend(doc.page_content.splitlines())
            return lines
        except ImportError:
            # 备选方案：使用 PyMuPDF
            try:
                import fitz
                doc = fitz.open(path)
                lines = []
                for page in doc:
                    lines.extend(page.get_text().splitlines())
                doc.close()
                return lines
            except ImportError:
                raise ImportError("pip install pdfplumber or pymupdf")
    
    def _extract_excel_lines(self, path: Path) -> List[str]:
        """提取 Excel 文档内容"""
        try:
            from langchain_community.document_loaders import UnstructuredExcelLoader
            loader = UnstructuredExcelLoader(str(path))
            docs = loader.load()
            lines = []
            for doc in docs:
                lines.extend(doc.page_content.splitlines())
            return lines
        except ImportError:
            # 备选方案：使用 openpyxl
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True)
                lines = []
                for sheet in wb.worksheets:
                    lines.append(f"=== Sheet: {sheet.title} ===")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = '\t'.join(str(c) if c is not None else '' for c in row)
                        if row_text.strip():
                            lines.append(row_text)
                wb.close()
                return lines
            except ImportError:
                raise ImportError("pip install openpyxl or unstructured")
    
    def _extract_ppt_lines(self, path: Path) -> List[str]:
        """提取 PowerPoint 文档内容"""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"=== Slide {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.extend(shape.text.splitlines())
            return lines
        except ImportError:
            raise ImportError("pip install python-pptx")
