"""
通用文档解析器 - 支持多种文档格式

从内存字节流解析文档，支持：
- Word: DOCX, DOC, RTF, ODT
- Excel: XLSX, XLS, CSV
- PowerPoint: PPTX, PPT
- PDF: PDF
- 文本: TXT, MD, HTML
- WPS: WPS, ET, DPS (通过转换)
"""

import io
import tempfile
import os
from typing import Dict, List, Optional, Tuple, Callable
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


class DocumentTypeDetector:
    """文档类型检测器 - 基于扩展名和 MIME 类型"""
    
    # 扩展名到 MIME 类型的映射
    EXTENSION_TO_MIME: Dict[str, str] = {
        # Word 文档
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.rtf': 'application/rtf',
        '.odt': 'application/vnd.oasis.opendocument.text',
        # Excel 表格
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.csv': 'text/csv',
        '.ods': 'application/vnd.oasis.opendocument.spreadsheet',
        # PowerPoint 演示
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.odp': 'application/vnd.oasis.opendocument.presentation',
        # PDF
        '.pdf': 'application/pdf',
        # 文本文件
        '.txt': 'text/plain',
        '.md': 'text/markdown',
        '.html': 'text/html',
        '.htm': 'text/html',
        # WPS Office (Kingsoft)
        '.wps': 'application/vnd.ms-works',  # WPS Writer
        '.et': 'application/vnd.ms-excel',   # WPS Spreadsheet (兼容)
        '.dps': 'application/vnd.ms-powerpoint',  # WPS Presentation (兼容)
    }
    
    # MIME 类型到文档类别的映射
    MIME_TO_CATEGORY: Dict[str, str] = {
        # Word 类
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'word',
        'application/msword': 'word',
        'application/rtf': 'word',
        'application/vnd.oasis.opendocument.text': 'word',
        # Excel 类
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'excel',
        'application/vnd.ms-excel': 'excel',
        'text/csv': 'excel',
        'application/vnd.oasis.opendocument.spreadsheet': 'excel',
        # PowerPoint 类
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'powerpoint',
        'application/vnd.ms-powerpoint': 'powerpoint',
        'application/vnd.oasis.opendocument.presentation': 'powerpoint',
        # PDF
        'application/pdf': 'pdf',
        # 文本类
        'text/plain': 'text',
        'text/markdown': 'text',
        'text/html': 'text',
        'text/csv': 'text',
    }
    
    @staticmethod
    def detect_type(filename: str, content_type: Optional[str] = None) -> Tuple[str, str]:
        """
        检测文档类型
        
        Args:
            filename: 文件名
            content_type: MIME 类型（可选）
            
        Returns:
            (文档类别, 扩展名) 元组
        """
        # 从文件名提取扩展名
        ext = Path(filename).suffix.lower()
        
        # 优先使用 content_type
        if content_type:
            category = DocumentTypeDetector.MIME_TO_CATEGORY.get(content_type)
            if category:
                return (category, ext)
        
        # 从扩展名推断
        if ext in DocumentTypeDetector.EXTENSION_TO_MIME:
            mime = DocumentTypeDetector.EXTENSION_TO_MIME[ext]
            category = DocumentTypeDetector.MIME_TO_CATEGORY.get(mime, 'unknown')
            return (category, ext)
        
        # 默认处理
        if ext in ['.txt', '.md', '.html', '.htm']:
            return ('text', ext)
        
        return ('unknown', ext)


class DocumentParser:
    """通用文档解析器 - 从内存字节流解析"""
    
    # 解析器映射：类别 -> 解析函数
    PARSERS: Dict[str, Callable[[bytes, str], str]] = {}
    
    @staticmethod
    def register_parser(category: str, parser_func: Callable[[bytes, str], str]):
        """注册解析器"""
        DocumentParser.PARSERS[category] = parser_func
    
    @staticmethod
    def parse(content_bytes: bytes, filename: str, content_type: Optional[str] = None) -> str:
        """
        解析文档内容
        
        Args:
            content_bytes: 文件字节流
            filename: 文件名
            content_type: MIME 类型（可选）
            
        Returns:
            解析后的文本内容
        """
        category, ext = DocumentTypeDetector.detect_type(filename, content_type)
        
        # 获取解析器
        parser = DocumentParser.PARSERS.get(category)
        if parser:
            try:
                return parser(content_bytes, ext)
            except Exception as e:
                logger.warning(f"解析器 {category} 失败，尝试通用方法: {e}")
        
        # 回退到通用文本解析
        return DocumentParser._parse_as_text(content_bytes, ext)
    
    @staticmethod
    def _parse_as_text(content_bytes: bytes, ext: str) -> str:
        """通用文本解析（UTF-8）"""
        try:
            return content_bytes.decode('utf-8', errors='replace')
        except Exception as e:
            raise ValueError(f"无法解析文件（扩展名: {ext}）: {e}")


# ===== Word 文档解析器 =====

def _parse_word_docx(content_bytes: bytes, ext: str) -> str:
    """解析 DOCX 文件（使用 python-docx）"""
    try:
        from docx import Document
        
        doc = Document(io.BytesIO(content_bytes))
        paragraphs = []
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 提取表格
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        
        return '\n'.join(paragraphs)
    except ImportError:
        # 回退到 Unstructured
        return _parse_word_unstructured(content_bytes, ext)


def _parse_word_unstructured(content_bytes: bytes, ext: str) -> str:
    """解析 Word 文件（使用 UnstructuredWordDocumentLoader）"""
    try:
        from langchain_community.document_loaders import UnstructuredWordDocumentLoader
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            tmp_file.write(content_bytes)
            tmp_path = tmp_file.name
        
        try:
            loader = UnstructuredWordDocumentLoader(tmp_path)
            documents = loader.load()
            return '\n\n'.join([doc.page_content for doc in documents])
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except ImportError:
        raise ImportError(
            "需要安装文档解析库。请运行：\n"
            "  pip install python-docx（推荐）\n"
            "  或 pip install unstructured[local-inference]"
        )


# ===== Excel 文档解析器 =====

def _parse_excel(content_bytes: bytes, ext: str) -> str:
    """解析 Excel 文件"""
    try:
        import pandas as pd
        
        # 根据扩展名选择引擎
        if ext == '.xlsx':
            engine = 'openpyxl'
        elif ext == '.xls':
            engine = 'xlrd'
        else:
            engine = None
        
        # 读取 Excel
        excel_file = io.BytesIO(content_bytes)
        df_dict = pd.read_excel(excel_file, engine=engine, sheet_name=None)
        
        # 转换为文本
        result = []
        for sheet_name, df in df_dict.items():
            result.append(f"=== 工作表: {sheet_name} ===")
            result.append(df.to_string(index=False))
            result.append("")  # 空行分隔
        
        return '\n'.join(result)
    except ImportError:
        # 回退到 Unstructured
        return _parse_excel_unstructured(content_bytes, ext)


def _parse_excel_unstructured(content_bytes: bytes, ext: str) -> str:
    """解析 Excel 文件（使用 UnstructuredExcelLoader）"""
    try:
        from langchain_community.document_loaders import UnstructuredExcelLoader
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            tmp_file.write(content_bytes)
            tmp_path = tmp_file.name
        
        try:
            loader = UnstructuredExcelLoader(tmp_path)
            documents = loader.load()
            return '\n\n'.join([doc.page_content for doc in documents])
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except ImportError:
        raise ImportError(
            "需要安装 Excel 解析库。请运行：\n"
            "  pip install pandas openpyxl xlrd（推荐）\n"
            "  或 pip install unstructured[local-inference]"
        )


def _parse_csv(content_bytes: bytes, ext: str) -> str:
    """解析 CSV 文件"""
    try:
        import pandas as pd
        
        df = pd.read_csv(io.BytesIO(content_bytes))
        return df.to_string(index=False)
    except ImportError:
        # 回退到文本解析
        return content_bytes.decode('utf-8', errors='replace')


# ===== PowerPoint 文档解析器 =====

def _parse_powerpoint(content_bytes: bytes, ext: str) -> str:
    """解析 PowerPoint 文件"""
    try:
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(content_bytes))
        result = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            result.append(f"=== 幻灯片 {slide_num} ===")
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result.append(shape.text)
            
            result.append("")  # 空行分隔
        
        return '\n'.join(result)
    except ImportError:
        raise ImportError(
            "需要安装 PowerPoint 解析库。请运行：\n"
            "  pip install python-pptx"
        )


# ===== PDF 文档解析器 =====

def _parse_pdf(content_bytes: bytes, ext: str) -> str:
    """解析 PDF 文件"""
    try:
        from langchain_community.document_loaders import PDFPlumberLoader
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(content_bytes)
            tmp_path = tmp_file.name
        
        try:
            loader = PDFPlumberLoader(tmp_path)
            documents = loader.load()
            return '\n\n'.join([doc.page_content for doc in documents])
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except ImportError:
        raise ImportError(
            "需要安装 PDF 解析库。请运行：\n"
            "  pip install pdfplumber 或 pip install pypdf"
        )


# ===== 文本文件解析器 =====

def _parse_text(content_bytes: bytes, ext: str) -> str:
    """解析文本文件"""
    # 尝试多种编码
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    
    for encoding in encodings:
        try:
            return content_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    
    # 如果都失败，使用 replace 模式
    return content_bytes.decode('utf-8', errors='replace')


# ===== 注册解析器 =====

DocumentParser.register_parser('word', _parse_word_docx)
DocumentParser.register_parser('excel', _parse_excel)
DocumentParser.register_parser('powerpoint', _parse_powerpoint)
DocumentParser.register_parser('pdf', _parse_pdf)
DocumentParser.register_parser('text', _parse_text)

# CSV 特殊处理
DocumentParser.register_parser('csv', _parse_csv)


__all__ = ['DocumentParser', 'DocumentTypeDetector']

